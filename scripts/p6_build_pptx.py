# -*- coding: utf-8 -*-
"""Phase 6: Abschlusspraesentation (python-pptx). 3 Sprecherteile, jede Inhaltsfolie mit Visual."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

ROOT=os.path.dirname(os.path.dirname(os.path.abspath(__file__))); CH=os.path.join(ROOT,"charts")
NAVY=RGBColor(0x1E,0x27,0x61); BLUE=RGBColor(0x00,0x72,0xB2); ORANGE=RGBColor(0xE6,0x9F,0x00)
DARKT=RGBColor(0x22,0x22,0x22); GRAY=RGBColor(0x66,0x66,0x66); WHITE=RGBColor(0xFF,0xFF,0xFF)
LIGHT=RGBColor(0xF2,0xF4,0xF8)
HEAD="Cambria"; BODY="Calibri"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
W,H=prs.slide_width,prs.slide_height

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(1,0,0,W,H); r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def tb(s,x,y,w,h,text,size,color=DARKT,bold=False,font=BODY,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,line=1.05):
    t=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=t.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    paras=text.split("\n")
    for i,ln in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=line
        r=p.add_run(); r.text=ln; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return t
def bullets(s,x,y,w,h,items,size=16,color=DARKT,gap=6):
    t=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=t.text_frame; tf.word_wrap=True
    for i,(txt,bold) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap); p.line_spacing=1.05
        r=p.add_run(); r.text=("•  "+txt); f=r.font; f.size=Pt(size); f.bold=bold; f.name=BODY; f.color.rgb=color
    return t
def chip(s,x,y,w,text,color):
    c=s.shapes.add_shape(5,Inches(x),Inches(y),Inches(w),Inches(0.42)); c.fill.solid(); c.fill.fore_color.rgb=color
    c.line.fill.background(); c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.size=Pt(12); r.font.bold=True; r.font.name=BODY; r.font.color.rgb=WHITE
def img(s,path,x,y,w=None,h=None):
    if os.path.exists(path): s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w) if w else None,Inches(h) if h else None)
def img_fit(s,path,bx,by,bw,bh):
    """Bettet ein Bild ein, in die Box (bx,by,bw,bh) eingepasst und zentriert (Seitenverhältnis erhalten)."""
    if not os.path.exists(path): return
    iw,ih=Image.open(path).size; a=iw/ih
    w=bw; h=w/a
    if h>bh: h=bh; w=h*a
    s.shapes.add_picture(path,Inches(bx+(bw-w)/2),Inches(by+(bh-h)/2),Inches(w),Inches(h))
def src(s,txt="Quelle: Destatis/Regionalstatistik (offene Daten), Abruf 2026-06-29 · eigene Berechnung"):
    tb(s,0.5,7.05,12.3,0.35,txt,9,GRAY)

# 1 TITEL (dunkel)
s=slide(NAVY)
tb(s,0.9,2.1,11.5,2.2,"Schulabschluss ist\nnicht nur Ländersache",46,WHITE,True,HEAD,line=1.02)
tb(s,0.95,4.5,11,0.5,"Wo Bildungsrisiken lokal sichtbar werden — eine Data Story mit offenen Daten",18,RGBColor(0xCA,0xDC,0xFC),font=BODY,italic=True)
tb(s,0.95,6.2,11,0.5,"Max Budde · John Kanto · Aaron Ziegler   |   HTW Berlin · W2-AA Analytische Anwendungen",13,RGBColor(0xCA,0xDC,0xFC))

# 2 VISION + Flow
s=slide()
tb(s,0.6,0.4,12,0.8,"Vision",34,NAVY,True,HEAD)
tb(s,0.6,1.45,7.1,2.6,"Bildungserfolg in Deutschland ist kein reines Länderphänomen — er wird lokal entschieden. Mit offenen Daten von Destatis und der Regionalstatistik verfolgen wir Schulabschlüsse von der Bundesland- bis auf die Kreisebene und verknüpfen sie mit Schulstruktur, Bildungsausgaben und Arbeitsmarkt.",17,DARKT,line=1.12)
# Flow chips
fx=0.7; fy=4.6
for i,(t,c) in enumerate([("INPUT\nAusgaben · Schulstruktur",BLUE),("OUTPUT\nAbschlüsse · Abgänge",RGBColor(0x6b,0x6e,0x2a)),("ÜBERGANG\nBerufsschule · Ausbildung",ORANGE),("ERGEBNIS\nEinkommen",RGBColor(0x55,0x52,0x77))]):
    box=s.shapes.add_shape(5,Inches(fx),Inches(fy),Inches(2.7),Inches(1.25)); box.fill.solid(); box.fill.fore_color.rgb=c; box.line.fill.background(); box.shadow.inherit=False
    tfx=box.text_frame; tfx.word_wrap=True; tfx.vertical_anchor=MSO_ANCHOR.MIDDLE
    for j,seg in enumerate(t.split("\n")):
        p=tfx.paragraphs[0] if j==0 else tfx.add_paragraph(); p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=seg; r.font.size=Pt(14 if j==0 else 10); r.font.bold=(j==0); r.font.color.rgb=WHITE; r.font.name=BODY
    if i<3: tb(s,fx+2.72,fy+0.35,0.5,0.5,"→",22,GRAY,True)
    fx+=3.08

# 3 LEITFRAGEN
s=slide()
tb(s,0.6,0.4,12,0.8,"Leitfragen",34,NAVY,True,HEAD)
cols=[("DER BEFUND",BLUE,["Welche Bundesländer führen 22/23 & 23/24 bei Abgängern ohne Abschluss?","Wo ist der Anteil ohne Hauptschulabschluss am höchsten?","Länder- oder Kreisproblem: Wie stark streuen die Kreise?"]),
      ("DIE STRUKTUR",ORANGE,["Schneiden Jungen und Mädchen unterschiedlich ab?","Wie prägt der Schulartmix die Abschlussverteilung?","Ändert sich die Wertung relativ statt absolut?"]),
      ("DAS ÖKONOMISCHE",RGBColor(0x55,0x52,0x77),["Wie verteilen sich die Bildungsausgaben nach Bereich?","Mehr Ausgaben je Schüler, bessere Abschlüsse?","Welche Kreise verbinden Bildungsrisiko, Arbeitslosigkeit & niedriges Einkommen?"])]
cx=0.6
for title,c,qs in cols:
    chip(s,cx,1.6,3.9,title,c)
    bullets(s,cx,2.25,3.9,4.5,[(q,False) for q in qs],size=13,gap=10)
    cx+=4.13

# Helper fuer LF-Inhaltsfolien
def lf_slide(part_chip, part_color, title, chart, findings, statbig=None, statlab=None):
    s=slide()
    chip(s,0.6,0.45,2.6,part_chip,part_color)
    tb(s,0.6,1.0,12,1.0,title,26,NAVY,True,HEAD)
    img_fit(s,os.path.join(CH,chart),6.55,1.9,6.5,5.0)
    bullets(s,0.6,2.1,5.9,4.3,[(f,b) for f,b in findings],size=15,gap=9)
    if statbig:
        tb(s,0.6,5.3,5.9,1.0,statbig,40,part_color,True,HEAD)
        tb(s,0.62,6.25,5.9,0.5,statlab,12,GRAY)
    src(s); return s

# TEIL 1 — BEFUND
lf_slide("TEIL 1 · BEFUND",BLUE,"LF1 — Sachsen-Anhalt führt bei Abgängen ohne Abschluss","pbi/pbi_lf1.png",
    [("Ost-Flächenländer + Bremen/SH an der Spitze",True),("Sachsen-Anhalt 2022: 11,3 % → 2023: 12,7 %",False),("Bayern am niedrigsten (5,4 %)",False),("Anteil bundesweit leicht steigend",False)],
    "12,7 %","Sachsen-Anhalt 2023 ohne Hauptschulabschluss")
lf_slide("TEIL 1 · BEFUND",BLUE,"LF2/LF3 — Es ist auch ein Kreisproblem","pbi/pbi_lf3.png",
    [("Hotspots: Anhalt-Bitterfeld 16,8 %, Pirmasens 16,5 %",True),("Starke Streuung INNERHALB der Länder",False),("Rheinland-Pfalz: Spannweite 12,7 Prozentpunkte",False),("Landesdurchschnitt verdeckt lokale Extreme",False)],
    "12,7 pp","größte Kreis-Spannweite (Rheinland-Pfalz)")

# TEIL 2 — STRUKTUR
lf_slide("TEIL 2 · STRUKTUR",ORANGE,"LF4 — Jungen öfter ohne Abschluss, Mädchen öfter Abitur","pbi/pbi_lf4.png",
    [("Ohne HSA: Jungen 8,4 % vs. Mädchen 5,8 %",True),("Abitur: Mädchen 37,1 % vs. Jungen 29,3 %",True),("Geschlechtergefälle in beide Richtungen",False),("Deutschland, Abgangsjahr 2023",False)],
    "+7,8 pp","Abitur-Vorsprung der Mädchen (DE 2023)")
# LF5 ohne eigenes Chart -> Stat-Layout
s=slide()
chip(s,0.6,0.45,2.6,"TEIL 2 · STRUKTUR",ORANGE)
tb(s,0.6,1.0,12,1.0,"LF5 — Der Schulartmix prägt die Abschlussverteilung",26,NAVY,True,HEAD)
mix=[("Förderschulen","42 %"),("Integr. Gesamtschulen","22 %"),("Schularten m. mehr. BG","16 %"),("Hauptschulen","13 %"),("Realschulen","5 %"),("Gymnasien (G8/G9)","~3 %")]
gx,gy=0.7,2.2
for i,(k,v) in enumerate(mix):
    bx=gx+(i%3)*4.1; by=gy+(i//3)*2.1
    card=s.shapes.add_shape(5,Inches(bx),Inches(by),Inches(3.8),Inches(1.8)); card.fill.solid(); card.fill.fore_color.rgb=LIGHT; card.line.fill.background(); card.shadow.inherit=False
    tb(s,bx+0.2,by+0.2,3.4,0.9,v,30,NAVY,True,HEAD)
    tb(s,bx+0.22,by+1.05,3.4,0.6,k,14,DARKT)
tb(s,0.6,6.2,12,0.5,"Herkunft der Abgänge ohne Hauptschulabschluss nach Schulart (DE 2023, Destatis 21111-12): Förderschulen stellen 42 %, integrierte Gesamtschulen 22 % – von Gymnasien nur rund 3 %.",13,GRAY,italic=True)
src(s)
lf_slide("TEIL 2 · STRUKTUR",ORANGE,"LF6 — Die Wertung kippt: absolut vs. relativ","pbi/pbi_lf6.png",
    [("Absolut führt NRW (11.835 ohne HSA)",False),("Relativ (je 1.000 der 15–18-Jährigen) führt Sachsen-Anhalt",True),("Kleine Ost-Länder & Bremen rücken nach vorn",False),("Relative Betrachtung ist fairer",True)],
    "Rang kippt","NRW #1 absolut → Sachsen-Anhalt #1 relativ")

# TEIL 3 — OEKONOMISCH
# LF7 stat layout
s=slide()
chip(s,0.6,0.45,2.9,"TEIL 3 · ÖKONOMISCH",RGBColor(0x55,0x52,0x77))
tb(s,0.6,1.0,12,1.0,"LF7 — Bildungsausgaben je Schüler nach Schulart",26,NAVY,True,HEAD)
au=[("Grundschulen","8.400 €"),("Realschulen","9.700 €"),("Gymnasien","10.900 €"),("Mehrere Bildungsg.","10.600 €"),("Integr. Gesamtschulen","11.600 €"),("Alle Schularten","9.800 €")]
for i,(k,v) in enumerate(au):
    bx=0.7+(i%3)*4.1; by=2.2+(i//3)*2.1
    card=s.shapes.add_shape(5,Inches(bx),Inches(by),Inches(3.8),Inches(1.8)); card.fill.solid(); card.fill.fore_color.rgb=LIGHT; card.line.fill.background(); card.shadow.inherit=False
    tb(s,bx+0.2,by+0.2,3.4,0.9,v,28,RGBColor(0x55,0x52,0x77),True,HEAD)
    tb(s,bx+0.22,by+1.05,3.4,0.6,k,14,DARKT)
tb(s,0.6,6.2,12,0.5,"Ausgaben je Schüler:in (Deutschland 2023). Gymnasien & Gesamtschulen am teuersten.",13,GRAY,italic=True)
src(s)
lf_slide("TEIL 3 · ÖKONOMISCH",RGBColor(0x55,0x52,0x77),"LF8 — Mehr Geld = mehr Abitur? Nur scheinbar","pbi/pbi_lf8.png",
    [("r = +0,61 über alle 16 BL (p=0,01) …",True),("… aber OHNE Stadtstaaten: r = −0,36 (n.s.)",True),("Positiver Eindruck = Stadtstaaten-Artefakt",True),("n=16 → kein robuster Befund, keine Kausalität",False)],
    "Artefakt","Stadtstaaten treiben die Trendlinie")
lf_slide("TEIL 3 · ÖKONOMISCH",RGBColor(0x55,0x52,0x77),"LF9 — Risiko-Kreise: Bildung trifft Arbeitsmarkt","pbi/pbi_lf9.png",
    [("Gelsenkirchen, Pirmasens, Mansfeld-Südharz führen",True),("Hohe Quote ohne HSA + hohe Jugend-ALQ + niedriges Einkommen",False),("Score z-standardisiert (3 Dimensionen), gleich gewichtet",False),("Datenstände 2023/2025/2021 dokumentiert",False)],
    "Top-Risiko","Gelsenkirchen: Score 8,1 · 13,0 % ohne HSA · 13,4 % Jugend-ALQ")

# DATENQUALITAET
s=slide()
tb(s,0.6,0.4,12,0.8,"Datenqualität & Methodik",30,NAVY,True,HEAD)
dq=[("6 offene Quellen","Destatis & Regionalstatistik, kein Login"),("Doppelte Absicherung","jede Kennzahl gegen amtliche Quellwerte nachgerechnet"),("Konsistenz","12/14 Länder: Kreis-Σ = Land exakt"),("Plausibilität","Σ Abschlussarten = Insgesamt (DE exakt)"),("Encoding","Windows-1252 → UTF-8 sauber"),("Geheimhaltung","–/./x als Missing, nie als 0")]
for i,(k,v) in enumerate(dq):
    bx=0.7+(i%3)*4.1; by=1.6+(i//3)*2.5
    card=s.shapes.add_shape(5,Inches(bx),Inches(by),Inches(3.8),Inches(2.2)); card.fill.solid(); card.fill.fore_color.rgb=LIGHT; card.line.fill.background(); card.shadow.inherit=False
    tb(s,bx+0.22,by+0.25,3.4,0.7,k,17,NAVY,True,HEAD)
    tb(s,bx+0.24,by+1.0,3.4,1.0,v,13,DARKT,line=1.1)
tb(s,0.6,6.6,12,0.5,"Sternschema: 4 Fakten (Abgänge, Schule, Arbeitsmarkt, Ausgaben) + 4 Dimensionen (Region, Zeit, Abschluss, Schulart).",12,GRAY,italic=True)

# FAZIT
s=slide(NAVY)
tb(s,0.9,1.0,11.5,0.8,"Fazit",34,WHITE,True,HEAD)
bullets(s,0.95,2.1,11.4,4.0,[
 ("Bundeslandpolitik prägt das Niveau — aber Bildungsrisiken entscheiden sich LOKAL.",True),
 ("Starke Kreis-Streuung: Landesdurchschnitte verdecken Hotspots.",False),
 ("Relative Betrachtung verändert die Rangfolge der Länder grundlegend.",False),
 ("Mehr Ausgaben gehen mit höheren Abiturquoten einher — kein Kausalbeweis.",False),
 ("Risiko-Kreise verbinden Bildungsrisiko und Jugendarbeitslosigkeit.",True),
],size=17,color=WHITE,gap=12)
tb(s,0.95,6.3,11.4,0.6,"Bildungsinvestitionen wirken nur, wenn sie lokal in Abschlüsse, Übergänge und Einkommen übersetzt werden.",14,RGBColor(0xCA,0xDC,0xFC),italic=True)

# DANKE
s=slide(NAVY)
tb(s,0.9,2.6,11.5,1.5,"Danke!",54,WHITE,True,HEAD)
tb(s,0.95,4.3,11,0.6,"Max Budde · John Kanto · Aaron Ziegler",18,RGBColor(0xCA,0xDC,0xFC))
tb(s,0.95,5.0,11,0.6,"HTW Berlin · W2-AA Analytische Anwendungen · Prof. Dr. Kempa",13,RGBColor(0xCA,0xDC,0xFC))

# ---- Sprechernotizen (3 Teile à ~5 Min) ----
NOTES=[
 "Sprecher: Max Budde (Teil 1 - Befund, ca. 5 Min). Begrüßung. Kernthese: Schulabschluss ist nicht nur Ländersache - Bildungsrisiken entscheiden sich LOKAL. Kurz Team (Budde/Kanto/Ziegler) und Modul W2-AA nennen. Überleitung zur Vision.",
 "Max Budde. Vision: Wir verfolgen Abschlüsse von der Bundesland- bis auf die Kreisebene und verknüpfen sie mit Schulstruktur, Ausgaben und Arbeitsmarkt. Daten-Flow am Schaubild durchgehen: INPUT -> OUTPUT -> ÜBERGANG -> ERGEBNIS. Betonen: ausschließlich offene Daten (Destatis/Regionalstatistik), kein Login.",
 "Max Budde. Die 9 Leitfragen in 3 Blöcken vorstellen: Befund (ich), Struktur (John), Ökonomie (Aaron). Nur kurz anreißen - Details folgen je Block.",
 "Max Budde - LF1. Sachsen-Anhalt führt 2023 mit 12,7 % Abgängen ohne Hauptschulabschluss; Ost-Flächenländer + Bremen/SH oben, Bayern am niedrigsten (5,4 %). Anstieg 2022 (11,3 %) -> 2023 (12,7 %). Wert in Power BI und unabhängig gegen die amtlichen Quellwerte verifiziert.",
 "Max Budde - LF2/LF3. Es ist auch ein Kreisproblem: Hotspots Anhalt-Bitterfeld 16,8 %, Pirmasens 16,5 %. Das Streudiagramm und die Streuungstabelle zeigen: starke Streuung INNERHALB der Länder (z. B. Rheinland-Pfalz Spannweite 12,7 pp). Landesdurchschnitt verdeckt lokale Extreme. Übergabe an John Kanto.",
 "Sprecher: John Kanto (Teil 2 - Struktur, ca. 5 Min). LF4 Geschlechtergefälle: ohne HSA Jungen 8,4 % vs. Mädchen 5,8 %; Abitur Mädchen 37,1 % vs. Jungen 29,3 %. Gefälle in beide Richtungen (DE 2023).",
 "John Kanto - LF5 Schulartmix (DE 2023): Der Schulartmix prägt die Abschlussverteilung. Von 55.705 Abgängen ohne Hauptschulabschluss stellen Förderschulen 42 % (23.324), integrierte Gesamtschulen 22 %, Schularten mit mehreren Bildungsgängen 16 % und Hauptschulen 13 % - von Gymnasien nur rund 3 % (Destatis 21111-12, Landesebene).",
 "John Kanto - LF6. Die Wertung kippt: absolut führt NRW (meiste Fälle), relativ (je 1.000 der 15-18-Jährigen) führt Sachsen-Anhalt. Kleine Ost-Länder und Bremen rücken nach vorn. Relative Betrachtung ist fairer. Übergabe an Aaron Ziegler.",
 "Sprecher: Aaron Ziegler (Teil 3 - Ökonomie, ca. 5 Min). LF7 Ausgaben je Schüler: Gymnasien/Gesamtschulen am teuersten (~11 Tsd €), Grundschulen ~8,4 Tsd €. Nach Bundesland führen die Stadtstaaten (Berlin/Hamburg).",
 "Aaron Ziegler - LF8. Mehr Ausgaben je Schüler gehen mit höherer Abiturquote einher: r = +0,61 über 16 Bundesländer; mit der Quote ohne HSA r = -0,33. WICHTIG: Zusammenhang, KEINE Kausalität; Strukturkosten der Stadtstaaten beachten.",
 "Aaron Ziegler - LF9. Risiko-Kreise verbinden hohes Bildungsrisiko (ohne HSA), hohe Jugend-Arbeitslosigkeit und niedriges verfügbares Einkommen: Gelsenkirchen führt (Score 8,1) vor Pirmasens und Mansfeld-Südharz. Score z-standardisiert über drei Dimensionen, gleich gewichtet. Datenstände 2023/2025/2021 transparent dokumentiert.",
 "Aaron Ziegler - Datenqualität & Methodik. 6 offene Quellen, jede Kennzahl unabhängig gegen die amtlichen Quellwerte nachgerechnet; Konsistenz/Plausibilität geprüft, Missing nie als 0. Behobene Datenfehler transparent dokumentiert (DQ8-DQ11): Dezimal-Locale (Jugend-ALQ x10), Whitespace-Join (Ausgaben), Jahresbezug 2023. Sternschema: 4 Kern-Fakten + 4 Dimensionen.",
 "Aaron Ziegler - Fazit. Kernbotschaften zusammenführen: Niveau = Länderpolitik, aber Risiken sind lokal; Streuung; relative Betrachtung verändert das Ranking; Ausgaben-Zusammenhang ohne Kausalbeweis; Risiko-Kreise. Schlusssatz: Bildungsinvestitionen wirken nur, wenn sie lokal in Abschlüsse, Übergänge und Einkommen übersetzt werden.",
 "Aaron Ziegler - Dank und Überleitung zur Diskussion. Für Rückfragen bereitstehen; Team nochmals nennen.",
]
for i,s in enumerate(prs.slides):
    if i<len(NOTES) and NOTES[i]:
        s.notes_slide.notes_text_frame.text=NOTES[i]

out=os.path.join(ROOT,"Schulabschluss_DataStory_Praesentation.pptx")
prs.save(out)
print("gespeichert:",out,"| Folien:",len(prs.slides._sldIdLst),"| Notizen:",sum(1 for n in NOTES if n))
